#!/usr/bin/env python3
"""
Document Processing API (FastAPI)
A FastAPI service for processing PDF, DOCX, and PPTX documents with Azure OpenAI-powered text extraction and analysis.
"""

import os
import tempfile
import json
import re
import time
from typing import Dict, List, Optional, Tuple
import logging
from contextlib import asynccontextmanager

# FastAPI imports
from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.responses import JSONResponse
from fastapi.middleware.cors import CORSMiddleware
import uvicorn

# Document processing imports
import PyPDF2
import pdfplumber
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from openai import AzureOpenAI
from dotenv import load_dotenv

# -----------------------------------------------------------------------------
# Logging
# -----------------------------------------------------------------------------
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("document-processor")

# -----------------------------------------------------------------------------
# Config
# -----------------------------------------------------------------------------
ALLOWED_EXTENSIONS = {"pdf", "docx", "pptx"}
MAX_CONTENT_LENGTH = 16 * 1024 * 1024  # 16MB
UPLOAD_FOLDER = tempfile.gettempdir()

# -----------------------------------------------------------------------------
# Helpers
# -----------------------------------------------------------------------------

def secure_filename(filename: str) -> str:
    """Very small, safe filename sanitizer suitable for APIs."""
    if not filename:
        return "upload"
    name = os.path.basename(filename)
    name = re.sub(r"[^A-Za-z0-9_.-]", "_", name)
    # Avoid dotfiles / empty names
    if name in {"", ".", ".."}:
        name = "upload"
    return name[:255]


def allowed_file(filename: str) -> bool:
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


async def save_upload_file_tmp(upload_file: UploadFile) -> str:
    """Save an UploadFile to a temp path with size check; return the path."""
    contents = await upload_file.read()
    if len(contents) > MAX_CONTENT_LENGTH:
        raise HTTPException(status_code=413, detail="File too large (max 16MB)")

    filename = secure_filename(upload_file.filename)
    file_path = os.path.join(UPLOAD_FOLDER, filename)
    with open(file_path, "wb") as f:
        f.write(contents)
    return file_path


# -----------------------------------------------------------------------------
# Core processor
# -----------------------------------------------------------------------------
class DocumentProcessor:
    """Advanced document processing with Azure OpenAI-powered skill extraction."""

    def __init__(self):
        load_dotenv()
        self.azure_openai_endpoint = os.getenv(
            "AZURE_OPENAI_ENDPOINT", "https://open-ai-res.openai.azure.com/"
        )
        self.azure_openai_key = os.getenv("AZURE_OPENAI_KEY")
        self.azure_openai_deployment = os.getenv("AZURE_OPENAI_DEPLOYMENT", "gpt-4o-mini")
        self.azure_openai_api_version = os.getenv("AZURE_OPENAI_API_VERSION", "2025-01-01-preview")

        logger.info(
            "Azure OpenAI Key: %s | Endpoint: %s | Deployment: %s | API: %s",
            "SET" if self.azure_openai_key else "NOT SET",
            self.azure_openai_endpoint,
            self.azure_openai_deployment,
            self.azure_openai_api_version,
        )

        self.azure_client: Optional[AzureOpenAI] = None
        if self.azure_openai_key:
            try:
                self.azure_client = AzureOpenAI(
                    azure_endpoint=self.azure_openai_endpoint,
                    api_key=self.azure_openai_key,
                    api_version=self.azure_openai_api_version,
                )
                logger.info("Azure OpenAI client initialized successfully")
            except Exception as e:
                logger.error("Failed to initialize Azure OpenAI client: %s", e)
        else:
            logger.warning("Azure OpenAI key not found, using fallback analysis")

    # ------------------------------- PDF ------------------------------------
    @staticmethod
    def extract_pdf_text_advanced(file_path: str) -> Tuple[str, Dict]:
        text_content = ""
        metadata = {"pageCount": 0, "method": "unknown", "hasImages": False}

        try:
            # Method 1: pdfplumber (best for complex layouts)
            try:
                with pdfplumber.open(file_path) as pdf:
                    pages_text = []
                    metadata["pageCount"] = len(pdf.pages)
                    for page_num, page in enumerate(pdf.pages, 1):
                        page_text = page.extract_text()
                        if page_text:
                            pages_text.append(f"--- Page {page_num} ---\n{page_text}")
                        if getattr(page, "images", None):
                            if page.images:
                                metadata["hasImages"] = True
                                pages_text.append(
                                    f"[Page {page_num} contains {len(page.images)} image(s)]"
                                )
                    if pages_text:
                        text_content = "\n\n".join(pages_text)
                        metadata["method"] = "pdfplumber"
                        logger.info(
                            "Extracted %d chars with pdfplumber", len(text_content)
                        )
            except Exception as e:
                logger.warning("pdfplumber failed: %s", e)

            # Method 2: PyMuPDF
            if len(text_content.strip()) < 100:
                try:
                    doc = fitz.open(file_path)
                    pages_text = []
                    metadata["pageCount"] = doc.page_count
                    for page_num in range(doc.page_count):
                        page = doc[page_num]
                        page_text = page.get_text()
                        if page_text.strip():
                            pages_text.append(f"--- Page {page_num + 1} ---\n{page_text}")
                        if page.get_images():
                            metadata["hasImages"] = True
                            pages_text.append(f"[Page {page_num + 1} contains images]")
                    doc.close()
                    if pages_text:
                        text_content = "\n\n".join(pages_text)
                        metadata["method"] = "pymupdf"
                        logger.info(
                            "Extracted %d chars with PyMuPDF", len(text_content)
                        )
                except Exception as e:
                    logger.warning("PyMuPDF failed: %s", e)

            # Method 3: PyPDF2
            if len(text_content.strip()) < 50:
                try:
                    with open(file_path, "rb") as file:
                        pdf_reader = PyPDF2.PdfReader(file)
                        pages_text = []
                        metadata["pageCount"] = len(pdf_reader.pages)
                        for page_num, page in enumerate(pdf_reader.pages, 1):
                            page_text = page.extract_text()
                            if page_text and page_text.strip():
                                pages_text.append(f"--- Page {page_num} ---\n{page_text}")
                        if pages_text:
                            text_content = "\n\n".join(pages_text)
                            metadata["method"] = "pypdf2"
                            logger.info(
                                "Extracted %d chars with PyPDF2", len(text_content)
                            )
                except Exception as e:
                    logger.error("PyPDF2 failed: %s", e)

            # Clean & finalize
            if text_content:
                text_content = DocumentProcessor.clean_text(text_content)
            else:
                text_content = (
                    "[PDF appears to be image-based or encrypted. Consider using OCR for better results.]"
                )
                metadata["method"] = "none"

        except Exception as e:
            logger.error("All PDF extraction methods failed: %s", e)
            text_content = f"[Error processing PDF: {str(e)}]"
            metadata["method"] = "error"

        return text_content, metadata

    # ------------------------------- DOCX -----------------------------------
    @staticmethod
    def extract_docx_text(file_path: str) -> Tuple[str, Dict]:
        try:
            doc = Document(file_path)
            paragraphs: List[str] = []
            for paragraph in doc.paragraphs:
                if paragraph.text and paragraph.text.strip():
                    paragraphs.append(paragraph.text)

            tables_text: List[str] = []
            for table in doc.tables:
                for row in table.rows:
                    row_text: List[str] = []
                    for cell in row.cells:
                        if cell.text and cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        tables_text.append(" | ".join(row_text))

            all_text = paragraphs + tables_text
            text_content = DocumentProcessor.clean_text("\n\n".join(all_text))
            metadata = {
                "paragraphCount": len(paragraphs),
                "tableCount": len(doc.tables),
                "method": "python-docx",
            }
            return text_content, metadata
        except Exception as e:
            logger.error("DOCX extraction failed: %s", e)
            return f"[Error processing DOCX: {str(e)}]", {"method": "error"}

    # ------------------------------- PPTX -----------------------------------
    @staticmethod
    def extract_pptx_text(file_path: str) -> Tuple[str, Dict]:
        try:
            prs = Presentation(file_path)
            slides_text: List[str] = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = [f"--- Slide {slide_num} ---"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        slide_content.append(shape.text.strip())
                # notes
                has_notes = getattr(slide, "has_notes_slide", False)
                if has_notes and getattr(slide, "notes_slide", None) and getattr(slide.notes_slide, "notes_text_frame", None):
                    notes_text = slide.notes_slide.notes_text_frame.text or ""
                    notes_text = notes_text.strip()
                    if notes_text:
                        slide_content.append(f"Notes: {notes_text}")
                if len(slide_content) > 1:
                    slides_text.append("\n".join(slide_content))

            text_content = DocumentProcessor.clean_text("\n\n".join(slides_text))
            metadata = {"slideCount": len(prs.slides), "method": "python-pptx"}
            return text_content, metadata
        except Exception as e:
            logger.error("PPTX extraction failed: %s", e)
            return f"[Error processing PPTX: {str(e)}]", {"method": "error"}

    # ------------------------------- Utils ----------------------------------
    @staticmethod
    def clean_text(text: str) -> str:
        text = re.sub(r"\n\s*\n", "\n\n", text)
        text = re.sub(r" +", " ", text)
        text = re.sub(r"^\s*\d+\s*$", "", text, flags=re.MULTILINE)
        text = re.sub(r"\s*\.\s*\.\s*\.+", "...", text)
        return text.strip()

    # ------------------------------- GenAI: Skills --------------------------
    def extract_skills_with_genai(self, text: str) -> Dict[str, List[str]]:
        if not self.azure_client:
            logger.error("Azure OpenAI client not available - cannot extract skills without GenAI")
            return self._empty_skills_result()
        try:
            prompt = self._create_skill_extraction_prompt(text)
            ai_response = self._call_azure_openai_for_skills(prompt)
            if ai_response:
                return self._parse_skills_response(ai_response)
            logger.error("No response from Azure OpenAI - skill extraction failed")
            return self._empty_skills_result()
        except Exception as e:
            logger.error("GenAI skill extraction failed: %s", e)
            return self._empty_skills_result()

    def _create_skill_extraction_prompt(self, text: str) -> str:
        max_text_length = 4000
        if len(text) > max_text_length:
            text = text[:max_text_length] + "... (truncated)"
        prompt = f"""
Analyze this document text and extract professional skills, competencies, and capabilities demonstrated by the person.

**Document Text:**
{text}

Please analyze this document thoroughly and provide a comprehensive JSON response with categorized skills:

1. **technical_skills**: Technical competencies including programming languages, frameworks, tools, platforms, databases, cloud services, software, methodologies
2. **programming_skills**: Specific programming languages and development technologies mentioned or demonstrated
3. **leadership_skills**: Leadership, management, mentoring, team coordination, strategic planning capabilities
4. **collaboration_skills**: Teamwork, cross-functional collaboration, stakeholder management, communication, partnership abilities
5. **research_skills**: Research methodology, data analysis, investigation, academic/scientific capabilities, analytical thinking
6. **soft_skills**: Problem-solving, creativity, adaptability, time management, critical thinking, emotional intelligence, communication

**Important Guidelines:**
- Extract skills that are explicitly mentioned OR clearly demonstrated through described activities/projects
- Include specific technologies, tools, programming languages, frameworks mentioned
- Identify skills from accomplishments and responsibilities described
- Be comprehensive but accurate - don't hallucinate skills not evidenced in the text
- Include both technical and soft skills
- Look for industry-specific knowledge and domain expertise
- Consider skills implied by job titles, projects, and achievements described

Please respond only with valid JSON in this exact format:
{{
  "technical_skills": ["skill1", "skill2", ...],
  "programming_skills": ["language1", "framework1", ...],
  "leadership_skills": ["skill1", "skill2", ...],
  "collaboration_skills": ["skill1", "skill2", ...],
  "research_skills": ["skill1", "skill2", ...],
  "soft_skills": ["skill1", "skill2", ...],
  "all_skills": ["combined list of all unique skills"]
}}
"""
        return prompt

    def _call_azure_openai_for_skills(self, prompt: str) -> Optional[str]:
        try:
            response = self.azure_client.chat.completions.create(
                model=self.azure_openai_deployment,
                messages=[
                    {
                        "role": "system",
                        "content": (
                            "You are an expert HR professional and technical recruiter specializing in skill assessment. "
                            "Analyze documents (resumes, project descriptions, reports, etc.) to extract professional skills and competencies. "
                            "Focus on both technical skills (programming, tools, platforms) and soft skills (leadership, communication, problem-solving). "
                            "Be thorough but accurate - only extract skills that are clearly evidenced in the text."
                        ),
                    },
                    {"role": "user", "content": prompt},
                ],
                max_tokens=2000,
                temperature=0.2,
                top_p=0.9,
            )
            if getattr(response, "choices", None):
                return response.choices[0].message.content
            return None
        except Exception as e:
            logger.error("Error calling Azure OpenAI for skills: %s", e)
            return None

    def _parse_skills_response(self, ai_response: str) -> Dict[str, List[str]]:
        try:
            json_start = ai_response.find("{")
            json_end = ai_response.rfind("}") + 1
            if json_start >= 0 and json_end > json_start:
                json_str = ai_response[json_start:json_end]
                parsed = json.loads(json_str)
                result = {
                    "technical_skills": parsed.get("technical_skills", []),
                    "programming_skills": parsed.get("programming_skills", []),
                    "leadership_skills": parsed.get("leadership_skills", []),
                    "collaboration_skills": parsed.get("collaboration_skills", []),
                    "research_skills": parsed.get("research_skills", []),
                    "soft_skills": parsed.get("soft_skills", []),
                    "all_skills": parsed.get("all_skills", []),
                }
                if not result["all_skills"]:
                    all_skills: List[str] = []
                    for k in [
                        "technical_skills",
                        "programming_skills",
                        "leadership_skills",
                        "collaboration_skills",
                        "research_skills",
                        "soft_skills",
                    ]:
                        all_skills.extend(result[k])
                    result["all_skills"] = sorted(set(all_skills))
                logger.info(
                    "GenAI extracted %d total skills", len(result.get("all_skills", []))
                )
                return result
            logger.error("Could not find valid JSON in AI response")
            return self._empty_skills_result()
        except json.JSONDecodeError as e:
            logger.error("JSON parsing error: %s", e)
            return self._empty_skills_result()
        except Exception as e:
            logger.error("Error parsing skills response: %s", e)
            return self._empty_skills_result()

    def _fallback_skill_extraction(self, text: str) -> Dict[str, List[str]]:
        logger.info("Using fallback pattern-based skill extraction")
        text_lower = text.lower()
        technical_patterns = [
            r"\b(?:python|javascript|java|react|typescript|node\.js|sql|aws|azure|docker|git)\b",
            r"\b(?:machine learning|data analysis|api|database|testing|devops|cloud)\b",
        ]
        soft_patterns = [
            r"(?:leadership|management|communication|collaboration|problem solving|teamwork)",
            r"(?:research|analysis|strategic planning|project management)",
        ]
        technical_skills: List[str] = []
        soft_skills: List[str] = []
        for pattern in technical_patterns:
            technical_skills.extend(re.findall(pattern, text_lower))
        for pattern in soft_patterns:
            soft_skills.extend(re.findall(pattern, text_lower))
        technical_skills = list({s.title() for s in technical_skills})
        soft_skills = list({s.title() for s in soft_skills})
        all_skills = technical_skills + soft_skills
        return {
            "technical_skills": technical_skills,
            "programming_skills": [
                s for s in technical_skills if s.lower() in ["python", "javascript", "java", "typescript"]
            ],
            "leadership_skills": [
                s for s in soft_skills if "leadership" in s.lower() or "management" in s.lower()
            ],
            "collaboration_skills": [
                s for s in soft_skills if "collaboration" in s.lower() or "teamwork" in s.lower()
            ],
            "research_skills": [
                s for s in soft_skills if "research" in s.lower() or "analysis" in s.lower()
            ],
            "soft_skills": soft_skills,
            "all_skills": all_skills,
        }

    def _empty_skills_result(self) -> Dict[str, List[str]]:
        return {
            "technical_skills": [],
            "programming_skills": [],
            "leadership_skills": [],
            "collaboration_skills": [],
            "research_skills": [],
            "soft_skills": [],
            "all_skills": [],
        }

    # ------------------------------- GenAI: Summary -------------------------
    def generate_summary_with_genai(self, text: str, categorized_skills: Dict[str, List[str]]) -> str:
        if not self.azure_client:
            logger.warning("Azure OpenAI client not available, falling back to empty summary")
            return ""
        try:
            prompt = self._create_summary_prompt(text, categorized_skills)
            ai_response = self._call_azure_openai_for_summary(prompt)
            if ai_response:
                return self._parse_summary_response(ai_response)
            logger.warning("No response from Azure OpenAI for summary")
            return ""
        except Exception as e:
            logger.error("GenAI summary generation failed: %s", e)
            return ""

    def _create_summary_prompt(self, text: str, categorized_skills: Dict[str, List[str]]) -> str:
        max_text_length = 3000
        if len(text) > max_text_length:
            text = text[:max_text_length] + "... (truncated)"
        skills_context = ""
        for category, skills in categorized_skills.items():
            if skills and category != "all_skills":
                skills_context += f"\n{category.replace('_', ' ').title()}: {', '.join(skills[:5])}"
        prompt = f"""
Analyze this document and the extracted skills to generate a comprehensive professional summary.

**Document Text:**
{text}

**Extracted Skills by Category:**
{skills_context}

Please generate a professional summary that:
1. Identifies the type of document (resume, project report, research paper, etc.)
2. Highlights the person's key professional strengths and expertise areas
3. Mentions specific technical skills and domains of experience
4. Describes the scope and depth of their capabilities
5. Provides quantitative context (experience level, project scale, etc.) when evident
6. Maintains a professional, objective tone

**Guidelines:**
- Write 2-3 sentences maximum
- Focus on the most significant skills and accomplishments
- Be specific about technical domains and expertise areas
- Mention leadership, collaboration, or research capabilities if evident
- Include context about the document's scope (comprehensive vs. focused, etc.)
- Don't make assumptions not supported by the evidence

Generate only the summary text, no additional formatting or explanation.
"""
        return prompt

    def _call_azure_openai_for_summary(self, prompt: str) -> Optional[str]:
        try:
            response = self.azure_client.chat.completions.create(
                model=self.azure_openai_deployment,
                messages=[
                    {
                        "role": "system",
                        "content": (
                            "You are an expert technical writer and HR professional specializing in creating concise, accurate professional summaries. "
                            "Generate clear, informative summaries that highlight key competencies and experience areas based on document analysis and extracted skills. "
                            "Focus on professional strengths, technical expertise, and demonstrated capabilities."
                        ),
                    },
                    {"role": "user", "content": prompt},
                ],
                max_tokens=300,
                temperature=0.3,
                top_p=0.9,
            )
            if getattr(response, "choices", None):
                return response.choices[0].message.content.strip()
            return None
        except Exception as e:
            logger.error("Error calling Azure OpenAI for summary: %s", e)
            return None

    def _parse_summary_response(self, ai_response: str) -> str:
        try:
            summary = ai_response.strip()
            summary = re.sub(r'^["\']|["\']$', "", summary)
            summary = re.sub(r"^\*\*|^##|^#", "", summary)
            summary = summary.strip()
            return summary or ""
        except Exception as e:
            logger.error("Error parsing summary response: %s", e)
            return ""

    # ------------------------------- GenAI: Suggested Skills ---------------
    def generate_suggested_skills_with_genai(self, text: str, existing_skills: List[str]) -> List[str]:
        if not self.azure_client:
            logger.warning("Azure OpenAI client not available, using fallback skill suggestions")
            return []
        try:
            prompt = self._create_skill_suggestions_prompt(text, existing_skills)
            ai_response = self._call_azure_openai_for_skill_suggestions(prompt)
            if ai_response:
                return self._parse_suggested_skills_response(ai_response)
            logger.warning("No response from Azure OpenAI for skill suggestions")
            return []
        except Exception as e:
            logger.error("GenAI skill suggestion failed: %s", e)
            return []

    def _create_skill_suggestions_prompt(self, text: str, existing_skills: List[str]) -> str:
        max_text_length = 3000
        if len(text) > max_text_length:
            text = text[:max_text_length] + "... (truncated)"
        existing_skills_str = ", ".join(existing_skills)
        prompt = f"""
Based on the document content and the skills already identified, suggest 10-15 additional relevant skills that this person likely possesses or should develop based on their field/domain.

**Document Text:**
{text}

**Already Identified Skills:**
{existing_skills_str}

**Instructions:**
- Suggest skills that are COMPLEMENTARY to what's already identified
- Focus on skills commonly needed in the same field/domain as the existing skills suggest
- Include both technical and soft skills that would be valuable
- Consider industry standards and common skill combinations
- Suggest skills that would enhance the person's profile
- Don't repeat skills already identified
- Prioritize practical, career-relevant skills

Please respond with exactly 10-15 skill suggestions in a simple JSON array format:
["skill1", "skill2", "skill3", ...]

Focus on high-value skills that would strengthen this professional profile.
"""
        return prompt

    def _call_azure_openai_for_skill_suggestions(self, prompt: str) -> Optional[str]:
        try:
            response = self.azure_client.chat.completions.create(
                model=self.azure_openai_deployment,
                messages=[
                    {
                        "role": "system",
                        "content": "You are an expert career advisor and skill analyzer. Provide helpful, relevant skill suggestions based on professional documents.",
                    },
                    {"role": "user", "content": prompt},
                ],
                max_tokens=500,
                temperature=0.3,
            )
            if getattr(response, "choices", None):
                return response.choices[0].message.content.strip()
            return None
        except Exception as e:
            logger.error("Azure OpenAI API call for skill suggestions failed: %s", e)
            return None

    def _parse_suggested_skills_response(self, response: str) -> List[str]:
        try:
            s = response.strip()
            if s.startswith("```json"):
                s = s[7:]
            if s.endswith("```"):
                s = s[:-3]
            s = s.strip()
            skills_list = json.loads(s)
            if isinstance(skills_list, list):
                out: List[str] = []
                for skill in skills_list:
                    if isinstance(skill, str) and skill.strip():
                        clean_skill = skill.strip().title()
                        if len(clean_skill) > 1 and clean_skill not in out:
                            out.append(clean_skill)
                return out[:15]
        except json.JSONDecodeError as e:
            logger.error("Failed to parse suggested skills JSON: %s", e)
        except Exception as e:
            logger.error("Error parsing suggested skills response: %s", e)
        return []

    # ------------------------------- GenAI: Title Generation ---------------
    def generate_title(self, text: str, filename: str) -> str:
        """Generate a concise title (4-5 words) from document text."""
        if not self.azure_client:
            logger.warning("Azure OpenAI client not available, falling back to filename-based title")
            return self._generate_title_from_filename(filename)
        
        try:
            # Truncate text for prompt
            max_text_length = 2000
            if len(text) > max_text_length:
                text = text[:max_text_length] + "... (truncated)"
            
            prompt = f"""
Generate a concise, professional title (4-5 words) for this document based on its content.

**Document Text:**
{text}

**Guidelines:**
- Create a title that accurately reflects the main topic or purpose
- Keep it concise (4-5 words maximum)
- Make it professional and clear
- Focus on the primary subject matter
- Avoid vague or generic titles
- Don't use phrases like "Document about..." or "Report on..."

Respond with ONLY the title text, nothing else.
"""
            response = self.azure_client.chat.completions.create(
                model=self.azure_openai_deployment,
                messages=[
                    {
                        "role": "system",
                        "content": "You are a professional editor who specializes in creating concise, accurate titles for documents."
                    },
                    {"role": "user", "content": prompt}
                ],
                max_tokens=20,
                temperature=0.3
            )
            
            if getattr(response, "choices", None):
                title = response.choices[0].message.content.strip()
                # Remove quotes if present
                title = re.sub(r'^["\']|["\']$', "", title)
                # Remove heading markdown if present
                title = re.sub(r"^#+\s*", "", title)
                return title or self._generate_title_from_filename(filename)
            return self._generate_title_from_filename(filename)
        except Exception as e:
            logger.error("Error generating title: %s", e)
            return self._generate_title_from_filename(filename)
    
    def _generate_title_from_filename(self, filename: str) -> str:
        """Generate a title from the filename as fallback."""
        try:
            # Remove extension
            base_name = os.path.splitext(filename)[0]
            # Replace underscores and hyphens with spaces
            clean_name = re.sub(r'[_-]', ' ', base_name)
            # Title case and limit length
            title = ' '.join(clean_name.split()[:5]).title()
            return title or "Uploaded Document"
        except Exception:
            return "Uploaded Document"
    
    # ------------------------------- GenAI: Description Generation ----------
    def generate_description(self, text: str) -> str:
        """Generate a concise description from document text."""
        if not self.azure_client:
            logger.warning("Azure OpenAI client not available, falling back to empty description")
            return ""
        
        try:
            # Truncate text for prompt
            max_text_length = 3000
            if len(text) > max_text_length:
                text = text[:max_text_length] + "... (truncated)"
            
            prompt = f"""
Generate a concise, professional description (1-2 sentences) for this document based on its content.

**Document Text:**
{text}

**Guidelines:**
- Create a description that summarizes the main topic or purpose
- Keep it concise (1-2 sentences, maximum 25 words)
- Make it professional and informative
- Focus on the key points or main subject matter
- Be specific rather than generic

Respond with ONLY the description text, nothing else.
"""
            response = self.azure_client.chat.completions.create(
                model=self.azure_openai_deployment,
                messages=[
                    {
                        "role": "system",
                        "content": "You are a professional editor who specializes in creating concise, accurate descriptions for documents."
                    },
                    {"role": "user", "content": prompt}
                ],
                max_tokens=60,
                temperature=0.3
            )
            
            if getattr(response, "choices", None):
                description = response.choices[0].message.content.strip()
                # Remove quotes if present
                description = re.sub(r'^["\']|["\']$', "", description)
                return description or ""
            return ""
        except Exception as e:
            logger.error("Error generating description: %s", e)
            return ""


# -----------------------------------------------------------------------------
# FastAPI app & lifespan
# -----------------------------------------------------------------------------
origins = [
    "http://localhost:5173",
    "http://localhost:3000",
    "http://localhost:8080",
    "http://localhost:8081",
    "http://127.0.0.1:5173",
    "http://127.0.0.1:8080",
    "http://127.0.0.1:8081",
]


@asynccontextmanager
async def lifespan(app: FastAPI):
    logger.info("Starting Document Processing API with Azure OpenAI integration (FastAPI)...")
    dp = DocumentProcessor()
    app.state.doc_processor = dp
    if dp.azure_client:
        logger.info("✅ Azure OpenAI client is ready for GenAI skill extraction")
    else:
        logger.warning("⚠️ Azure OpenAI client not available - using fallback pattern-based extraction where applicable")
    yield
    # No teardown needed


app = FastAPI(
    title="Document Processing API",
    version="1.0.0",
    lifespan=lifespan,
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# -----------------------------------------------------------------------------
# Routes
# -----------------------------------------------------------------------------
@app.get("/health")
async def health_check():
    return {"status": "healthy", "service": "document-processor-api"}


@app.post("/process-document")
async def process_document(file: UploadFile = File(...)):
    start_time = time.time()

    if not file or not file.filename:
        raise HTTPException(status_code=400, detail="No file provided")

    if not allowed_file(file.filename):
        raise HTTPException(
            status_code=400,
            detail=f"File type not supported. Allowed types: {', '.join(sorted(ALLOWED_EXTENSIONS)).upper()}",
        )

    file_path = await save_upload_file_tmp(file)

    try:
        filename = os.path.basename(file_path)
        file_extension = filename.lower().split(".")[-1]

        if file_extension == "pdf":
            text, metadata = DocumentProcessor.extract_pdf_text_advanced(file_path)
        elif file_extension == "docx":
            text, metadata = DocumentProcessor.extract_docx_text(file_path)
        elif file_extension == "pptx":
            text, metadata = DocumentProcessor.extract_pptx_text(file_path)
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type")

        dp: DocumentProcessor = app.state.doc_processor

        # Extract skills and generate summary using GenAI (if configured)
        categorized_skills = dp.extract_skills_with_genai(text)
        summary = dp.generate_summary_with_genai(text, categorized_skills)
        suggested_skills = dp.generate_suggested_skills_with_genai(
            text, categorized_skills.get("all_skills", [])
        )
        
        # Generate title and description
        title = dp.generate_title(text, filename)
        description = dp.generate_description(text)

        processing_time = int((time.time() - start_time) * 1000)

        response_data = {
            "text": text,
            "skills": categorized_skills.get("all_skills", []),
            "categorized_skills": categorized_skills,
            "suggested_skills": suggested_skills,
            "summary": summary,
            "generated_title": title,
            "generated_description": description,
            "metadata": {
                **metadata,
                "wordCount": len(text.split()) if isinstance(text, str) else 0,
                "fileType": file_extension.upper(),
                "processingTime": processing_time,
                "fileName": filename,
                "skillCategories": {
                    "technical": len(categorized_skills.get("technical_skills", [])),
                    "programming": len(categorized_skills.get("programming_skills", [])),
                    "leadership": len(categorized_skills.get("leadership_skills", [])),
                    "collaboration": len(categorized_skills.get("collaboration_skills", [])),
                    "research": len(categorized_skills.get("research_skills", [])),
                    "soft": len(categorized_skills.get("soft_skills", [])),
                },
            },
        }
        logger.info(
            "Processed %s: %d chars, %d skills",
            filename,
            len(text) if isinstance(text, str) else 0,
            len(response_data["skills"]),
        )
        return JSONResponse(content=response_data)

    except HTTPException:
        raise
    except Exception as e:
        logger.exception("Error processing document: %s", e)
        raise HTTPException(status_code=500, detail=f"Processing failed: {str(e)}")
    finally:
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
        except Exception:
            pass


@app.post("/extract-text")
async def extract_text_only(file: UploadFile = File(...)):
    if not file or not file.filename:
        raise HTTPException(status_code=400, detail="No file provided")

    if not allowed_file(file.filename):
        raise HTTPException(status_code=400, detail="File type not supported")

    file_path = await save_upload_file_tmp(file)

    try:
        file_extension = file.filename.lower().split(".")[-1]
        if file_extension == "pdf":
            text, metadata = DocumentProcessor.extract_pdf_text_advanced(file_path)
        elif file_extension == "docx":
            text, metadata = DocumentProcessor.extract_docx_text(file_path)
        elif file_extension == "pptx":
            text, metadata = DocumentProcessor.extract_pptx_text(file_path)
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type")

        return JSONResponse(content={"text": text, "metadata": metadata})

    except HTTPException:
        raise
    except Exception as e:
        logger.exception("Error extracting text: %s", e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
        except Exception:
            pass


@app.post("/generate-metadata")
async def generate_metadata(file: UploadFile = File(...)):
    """Generate a concise title and description from an uploaded document."""
    start_time = time.time()

    if not file or not file.filename:
        raise HTTPException(status_code=400, detail="No file provided")

    if not allowed_file(file.filename):
        raise HTTPException(
            status_code=400,
            detail=f"File type not supported. Allowed types: {', '.join(sorted(ALLOWED_EXTENSIONS)).upper()}",
        )

    file_path = await save_upload_file_tmp(file)

    try:
        filename = os.path.basename(file_path)
        file_extension = filename.lower().split(".")[-1]

        # Extract text from document
        if file_extension == "pdf":
            text, metadata = DocumentProcessor.extract_pdf_text_advanced(file_path)
        elif file_extension == "docx":
            text, metadata = DocumentProcessor.extract_docx_text(file_path)
        elif file_extension == "pptx":
            text, metadata = DocumentProcessor.extract_pptx_text(file_path)
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type")

        dp: DocumentProcessor = app.state.doc_processor
        title = dp.generate_title(text, filename)
        description = dp.generate_description(text)

        processing_time = int((time.time() - start_time) * 1000)

        response_data = {
            "title": title,
            "description": description,
            "metadata": {
                **metadata,
                "wordCount": len(text.split()) if isinstance(text, str) else 0,
                "fileType": file_extension.upper(),
                "processingTime": processing_time,
                "fileName": filename
            }
        }
        
        return JSONResponse(content=response_data)

    except HTTPException:
        raise
    except Exception as e:
        logger.exception("Error generating metadata: %s", e)
        raise HTTPException(status_code=500, detail=f"Metadata generation failed: {str(e)}")
    finally:
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
        except Exception:
            pass


# -----------------------------------------------------------------------------
# Entrypoint
# -----------------------------------------------------------------------------
if __name__ == "__main__":
    logger.info("Allowed file types: %s", ", ".join(sorted(ALLOWED_EXTENSIONS)))
    logger.info("Max file size: %.1fMB", MAX_CONTENT_LENGTH / (1024 * 1024))
    uvicorn.run(
        "app:app",
        host="0.0.0.0",
        port=5001,
        reload=True,
    )
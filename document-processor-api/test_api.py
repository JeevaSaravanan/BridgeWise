#!/usr/bin/env python3
"""
Test script for the Document Processing API
"""

import requests
import json
import os
import tempfile
from docx import Document
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter

# Test the API health
def test_health():
    print("Testing API health...")
    try:
        response = requests.get('http://localhost:5001/health')
        print(f"Health check: {response.status_code} - {response.json()}")
        return response.status_code == 200
    except Exception as e:
        print(f"Health check failed: {e}")
        return False

def create_test_pdf():
    """Create a test PDF file"""
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
    
    c = canvas.Canvas(temp_file.name, pagesize=letter)
    c.drawString(72, 750, "John Doe - Software Engineer")
    c.drawString(72, 720, "")
    c.drawString(72, 690, "Experience:")
    c.drawString(72, 660, "• 5 years of Python development")
    c.drawString(72, 630, "• React and JavaScript expertise")
    c.drawString(72, 600, "• AWS cloud architecture")
    c.drawString(72, 570, "• Machine learning projects")
    c.drawString(72, 540, "• Team leadership and project management")
    c.drawString(72, 510, "")
    c.drawString(72, 480, "Skills:")
    c.drawString(72, 450, "Python, React, JavaScript, TypeScript, AWS, Docker,")
    c.drawString(72, 420, "Kubernetes, Machine Learning, Data Analysis,")
    c.drawString(72, 390, "PostgreSQL, MongoDB, Git, Agile")
    c.drawString(72, 360, "")
    c.drawString(72, 330, "Education:")
    c.drawString(72, 300, "Bachelor of Computer Science")
    c.save()
    
    return temp_file.name

def create_test_docx():
    """Create a test DOCX file"""
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.docx')
    
    doc = Document()
    doc.add_heading('John Doe - Software Engineer', 0)
    
    doc.add_heading('Experience', level=1)
    doc.add_paragraph('• 5 years of Python development')
    doc.add_paragraph('• React and JavaScript expertise')
    doc.add_paragraph('• AWS cloud architecture')
    doc.add_paragraph('• Machine learning projects')
    doc.add_paragraph('• Team leadership and project management')
    
    doc.add_heading('Skills', level=1)
    doc.add_paragraph('Python, React, JavaScript, TypeScript, AWS, Docker, Kubernetes, Machine Learning, Data Analysis, PostgreSQL, MongoDB, Git, Agile')
    
    doc.add_heading('Education', level=1)
    doc.add_paragraph('Bachelor of Computer Science')
    
    doc.save(temp_file.name)
    return temp_file.name

def create_test_pptx():
    """Create a test PPTX file"""
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    
    prs = Presentation()
    
    # Slide 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "John Doe - Software Engineer"
    slide1.shapes.placeholders[1].text = "Professional Portfolio Overview"
    
    # Slide 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Experience"
    slide2.shapes.placeholders[1].text = """• 5 years of Python development
• React and JavaScript expertise
• AWS cloud architecture
• Machine learning projects
• Team leadership and project management"""
    
    # Slide 3
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Skills"
    slide3.shapes.placeholders[1].text = """Python, React, JavaScript, TypeScript, AWS, Docker, Kubernetes, Machine Learning, Data Analysis, PostgreSQL, MongoDB, Git, Agile"""
    
    prs.save(temp_file.name)
    return temp_file.name

def test_file_upload(file_path, file_type):
    """Test uploading a file to the API"""
    print(f"\nTesting {file_type} upload: {os.path.basename(file_path)}")
    
    try:
        with open(file_path, 'rb') as f:
            files = {'file': (os.path.basename(file_path), f, f'application/{file_type.lower()}')}
            
            # Test full processing endpoint
            print("Testing /process-document endpoint...")
            response = requests.post('http://localhost:5001/process-document', files=files)
            
            if response.status_code == 200:
                result = response.json()
                print(f"✓ Success! Extracted {len(result['text'])} characters")
                print(f"✓ Found {len(result['skills'])} skills: {', '.join(result['skills'][:5])}{'...' if len(result['skills']) > 5 else ''}")
                print(f"✓ Generated summary: {result['summary'][:100]}...")
                print(f"✓ Processing time: {result['metadata']['processingTime']}ms")
                print(f"✓ Method used: {result['metadata'].get('method', 'unknown')}")
                
                return True
            else:
                print(f"✗ Failed: {response.status_code} - {response.text}")
                return False
                
    except Exception as e:
        print(f"✗ Upload test failed: {e}")
        return False

def test_specific_pdf(pdf_path):
    """Test with a specific PDF file provided by user"""
    print(f"\nTesting specific PDF file: {pdf_path}")
    
    if not os.path.exists(pdf_path):
        print(f"✗ File not found: {pdf_path}")
        return False
    
    try:
        file_size = os.path.getsize(pdf_path) / (1024 * 1024)  # Size in MB
        print(f"📁 File size: {file_size:.2f} MB")
        
        if file_size > 16:  # API has 16MB limit
            print("⚠️  Warning: File exceeds 16MB limit, may be rejected by API")
        
        with open(pdf_path, 'rb') as f:
            files = {'file': (os.path.basename(pdf_path), f, 'application/pdf')}
            
            print("🔄 Processing document with Flask API...")
            response = requests.post('http://localhost:5001/process-document', files=files)
            
            if response.status_code == 200:
                result = response.json()
                print(f"✅ SUCCESS! Document processed")
                print(f"📄 Extracted text: {len(result['text'])} characters")
                print(f"🎯 Found skills ({len(result['skills'])}): {', '.join(result['skills'][:8])}{'...' if len(result['skills']) > 8 else ''}")
                print(f"📝 Summary: {result['summary']}")
                print(f"⏱️  Processing time: {result['metadata']['processingTime']}ms")
                print(f"🔧 Extraction method: {result['metadata'].get('method', 'unknown')}")
                print(f"📊 Word count: {result['metadata']['wordCount']}")
                print(f"📑 Page count: {result['metadata'].get('pageCount', 'unknown')}")
                
                # Show first 500 characters of extracted text
                print(f"\n📄 Text preview (first 500 chars):")
                print("-" * 50)
                print(result['text'][:500] + "..." if len(result['text']) > 500 else result['text'])
                print("-" * 50)
                
                return True
            else:
                print(f"❌ FAILED: {response.status_code}")
                try:
                    error_detail = response.json()
                    print(f"Error details: {error_detail}")
                except:
                    print(f"Error response: {response.text}")
                return False
                
    except Exception as e:
        print(f"❌ Test failed with exception: {e}")
        return False

def test_document_processing():
    """Test document processing with actual file uploads"""
    print("\nTesting document processing with real files...")
    
    # Test specific user PDF first
    user_pdf = "/Users/rohan/Documents/MSU/Classes/MTH 890 Capstone/February progress review Stunio.pdf"
    
    if os.path.exists(user_pdf):
        print("🎯 Testing with user's specific PDF file...")
        test_specific_pdf(user_pdf)
    else:
        print(f"⚠️  User PDF not found at: {user_pdf}")
        print("Will create and test with sample files instead...\n")
    
    test_files = []
    
    try:
        # Create test files
        print("Creating sample test files...")
        
        # Create PDF
        try:
            pdf_file = create_test_pdf()
            test_files.append((pdf_file, 'PDF'))
            print(f"✓ Created test PDF: {os.path.basename(pdf_file)}")
        except Exception as e:
            print(f"✗ Failed to create PDF: {e}")
        
        # Create DOCX
        try:
            docx_file = create_test_docx()
            test_files.append((docx_file, 'DOCX'))
            print(f"✓ Created test DOCX: {os.path.basename(docx_file)}")
        except Exception as e:
            print(f"✗ Failed to create DOCX: {e}")
        
        # Create PPTX
        try:
            pptx_file = create_test_pptx()
            test_files.append((pptx_file, 'PPTX'))
            print(f"✓ Created test PPTX: {os.path.basename(pptx_file)}")
        except Exception as e:
            print(f"✗ Failed to create PPTX: {e}")
        
        # Test each file
        success_count = 0
        for file_path, file_type in test_files:
            if test_file_upload(file_path, file_type):
                success_count += 1
        
        print(f"\n📊 Sample Files Test Results: {success_count}/{len(test_files)} files processed successfully")
        
    except requests.exceptions.ConnectionError:
        print("ERROR: Cannot connect to Flask API. Make sure it's running on port 5001.")
    except Exception as e:
        print(f"Test failed: {e}")
    finally:
        # Clean up test files
        print("\nCleaning up sample test files...")
        for file_path, _ in test_files:
            if os.path.exists(file_path):
                os.remove(file_path)
                print(f"✓ Removed {os.path.basename(file_path)}")

if __name__ == "__main__":
    print("Document Processing API Test")
    print("=" * 40)
    test_health()
    test_document_processing()
    print("\n" + "=" * 40)
    print("Test completed!")
    print("\nTo test with real documents:")
    print("1. Open the web interface at http://localhost:8081")
    print("2. Click 'Create Portfolio Item'")
    print("3. Upload a PDF, DOCX, or PPTX file")
    print("4. The Flask API will process it and extract text, skills, and summary")

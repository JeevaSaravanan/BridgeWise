#!/usr/bin/env python3
"""
Document Processing API
A Flask API for processing PDF, DOCX, and PPTX documents with Azure OpenAI-powered text extraction and analysis.
"""

import os
import io
import tempfile
import json
from flask import Flask, request, jsonify
from flask_cors import CORS
import PyPDF2
import pdfplumber
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import re
import time
from typing import Dict, List, Optional, Tuple
from werkzeug.utils import secure_filename
import logging
from openai import AzureOpenAI

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__)
CORS(app, origins=[
    "http://localhost:5173", 
    "http://localhost:3000", 
    "http://localhost:8080",
    "http://localhost:8081",
    "http://127.0.0.1:5173",
    "http://127.0.0.1:8080",
    "http://127.0.0.1:8081"
])

# Configuration
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size
ALLOWED_EXTENSIONS = {'pdf', 'docx', 'pptx'}
UPLOAD_FOLDER = tempfile.gettempdir()

def allowed_file(filename: str) -> bool:
    """Check if file extension is allowed."""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

class DocumentProcessor:
    """Advanced document processing with Azure OpenAI-powered skill extraction."""
    
    def __init__(self):
        """Initialize DocumentProcessor with Azure OpenAI configuration."""
        # Azure OpenAI configuration
        self.azure_openai_endpoint = os.getenv("AZURE_OPENAI_ENDPOINT")
        self.azure_openai_key = os.getenv("AZURE_OPENAI_API_KEY")
        self.azure_openai_deployment = os.getenv("AZURE_OPENAI_DEPLOYMENT_NAME", "gpt-4")
        
        # Initialize Azure OpenAI client
        self.azure_client = None
        if self.azure_openai_endpoint and self.azure_openai_key:
            try:
                self.azure_client = AzureOpenAI(
                    azure_endpoint=self.azure_openai_endpoint,
                    api_key=self.azure_openai_key,
                    api_version="2024-02-01"
                )
                logger.info("Azure OpenAI client initialized successfully")
            except Exception as e:
                logger.error(f"Failed to initialize Azure OpenAI client: {e}")
                self.azure_client = None
        else:
            logger.warning("Azure OpenAI credentials not found in environment variables")
    
    @staticmethod
    def extract_pdf_text_advanced(file_path: str) -> Tuple[str, Dict]:
        """Extract text from PDF using multiple methods for best results."""
        text_content = ""
        metadata = {"pageCount": 0, "method": "unknown", "hasImages": False}
        
        try:
            # Method 1: Try pdfplumber first (best for complex layouts)
            try:
                with pdfplumber.open(file_path) as pdf:
                    pages_text = []
                    metadata["pageCount"] = len(pdf.pages)
                    
                    for page_num, page in enumerate(pdf.pages, 1):
                        page_text = page.extract_text()
                        if page_text:
                            pages_text.append(f"--- Page {page_num} ---\n{page_text}")
                        
                        # Check for images/charts
                        if page.images:
                            metadata["hasImages"] = True
                            pages_text.append(f"[Page {page_num} contains {len(page.images)} image(s)]")
                    
                    if pages_text:
                        text_content = "\n\n".join(pages_text)
                        metadata["method"] = "pdfplumber"
                        logger.info(f"Successfully extracted text using pdfplumber: {len(text_content)} chars")
            except Exception as e:
                logger.warning(f"pdfplumber failed: {e}")
            
            # Method 2: Fallback to PyMuPDF if pdfplumber fails or returns little text
            if len(text_content.strip()) < 100:
                try:
                    doc = fitz.open(file_path)
                    pages_text = []
                    metadata["pageCount"] = doc.page_count
                    
                    for page_num in range(doc.page_count):
                        page = doc[page_num]
                        page_text = page.get_text()
                        if page_text.strip():
                            pages_text.append(f"--- Page {page_num + 1} ---\n{page_text}")
                        
                        # Check for images
                        if page.get_images():
                            metadata["hasImages"] = True
                            pages_text.append(f"[Page {page_num + 1} contains images]")
                    
                    doc.close()
                    
                    if pages_text:
                        text_content = "\n\n".join(pages_text)
                        metadata["method"] = "pymupdf"
                        logger.info(f"Successfully extracted text using PyMuPDF: {len(text_content)} chars")
                except Exception as e:
                    logger.warning(f"PyMuPDF failed: {e}")
            
            # Method 3: Final fallback to PyPDF2
            if len(text_content.strip()) < 50:
                try:
                    with open(file_path, 'rb') as file:
                        pdf_reader = PyPDF2.PdfReader(file)
                        pages_text = []
                        metadata["pageCount"] = len(pdf_reader.pages)
                        
                        for page_num, page in enumerate(pdf_reader.pages, 1):
                            page_text = page.extract_text()
                            if page_text.strip():
                                pages_text.append(f"--- Page {page_num} ---\n{page_text}")
                        
                        if pages_text:
                            text_content = "\n\n".join(pages_text)
                            metadata["method"] = "pypdf2"
                            logger.info(f"Successfully extracted text using PyPDF2: {len(text_content)} chars")
                except Exception as e:
                    logger.error(f"PyPDF2 failed: {e}")
            
            # Clean up the extracted text
            if text_content:
                text_content = DocumentProcessor.clean_text(text_content)
            else:
                text_content = "[PDF appears to be image-based or encrypted. Consider using OCR for better results.]"
                metadata["method"] = "none"
            
        except Exception as e:
            logger.error(f"All PDF extraction methods failed: {e}")
            text_content = f"[Error processing PDF: {str(e)}]"
            metadata["method"] = "error"
        
        return text_content, metadata
    
    @staticmethod
    def extract_docx_text(file_path: str) -> Tuple[str, Dict]:
        """Extract text from DOCX file."""
        try:
            doc = Document(file_path)
            paragraphs = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    paragraphs.append(paragraph.text)
            
            # Extract text from tables
            tables_text = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        tables_text.append(" | ".join(row_text))
            
            all_text = paragraphs + tables_text
            text_content = "\n\n".join(all_text)
            text_content = DocumentProcessor.clean_text(text_content)
            
            metadata = {
                "paragraphCount": len(paragraphs),
                "tableCount": len(doc.tables),
                "method": "python-docx"
            }
            
            return text_content, metadata
            
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            return f"[Error processing DOCX: {str(e)}]", {"method": "error"}
    
    @staticmethod
    def extract_pptx_text(file_path: str) -> Tuple[str, Dict]:
        """Extract text from PPTX file."""
        try:
            prs = Presentation(file_path)
            slides_text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = [f"--- Slide {slide_num} ---"]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                
                # Extract notes
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_content.append(f"Notes: {notes_text}")
                
                if len(slide_content) > 1:  # More than just the slide header
                    slides_text.append("\n".join(slide_content))
            
            text_content = "\n\n".join(slides_text)
            text_content = DocumentProcessor.clean_text(text_content)
            
            metadata = {
                "slideCount": len(prs.slides),
                "method": "python-pptx"
            }
            
            return text_content, metadata
            
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            return f"[Error processing PPTX: {str(e)}]", {"method": "error"}
    
    @staticmethod
    def clean_text(text: str) -> str:
        """Clean and normalize extracted text."""
        # Remove excessive whitespace
        text = re.sub(r'\n\s*\n', '\n\n', text)
        text = re.sub(r' +', ' ', text)
        
        # Remove common PDF artifacts
        text = re.sub(r'^\s*\d+\s*$', '', text, flags=re.MULTILINE)  # Page numbers on their own lines
        text = re.sub(r'\s*\.\s*\.\s*\.+', '...', text)  # Multiple dots
        
        return text.strip()
    
    def extract_skills_with_genai(self, text: str) -> Dict[str, List[str]]:
        """Extract skills from text using Azure OpenAI GenAI analysis."""
        if not self.azure_client:
            logger.warning("Azure OpenAI client not available, falling back to pattern-based analysis")
            return self._fallback_skill_extraction(text)
        
        try:
            # Create comprehensive prompt for skill extraction
            prompt = self._create_skill_extraction_prompt(text)
            
            # Call Azure OpenAI
            ai_response = self._call_azure_openai_for_skills(prompt)
            
            if ai_response:
                return self._parse_skills_response(ai_response)
            else:
                logger.warning("No response from Azure OpenAI, using fallback")
                return self._fallback_skill_extraction(text)
                
        except Exception as e:
            logger.error(f"GenAI skill extraction failed: {e}")
            return self._fallback_skill_extraction(text)
    
    def _create_skill_extraction_prompt(self, text: str) -> str:
        """Create a comprehensive prompt for skill extraction from document text."""
        
        # Truncate text if too long to fit in prompt
        max_text_length = 4000
        if len(text) > max_text_length:
            text = text[:max_text_length] + "... (truncated)"
        
        prompt = f"""
Analyze this document text and extract professional skills, competencies, and capabilities demonstrated by the person.

**Document Text:**
{text}

Please analyze this document thoroughly and provide a comprehensive JSON response with categorized skills:

1. **technical_skills**: Technical competencies including programming languages, frameworks, tools, platforms, databases, cloud services, software, methodologies
2. **programming_skills**: Specific programming languages and development technologies mentioned or demonstrated
3. **leadership_skills**: Leadership, management, mentoring, team coordination, strategic planning capabilities
4. **collaboration_skills**: Teamwork, cross-functional collaboration, stakeholder management, communication, partnership abilities
5. **research_skills**: Research methodology, data analysis, investigation, academic/scientific capabilities, analytical thinking
6. **soft_skills**: Problem-solving, creativity, adaptability, time management, critical thinking, emotional intelligence, communication

**Important Guidelines:**
- Extract skills that are explicitly mentioned OR clearly demonstrated through described activities/projects
- Include specific technologies, tools, programming languages, frameworks mentioned
- Identify skills from accomplishments and responsibilities described
- Be comprehensive but accurate - don't hallucinate skills not evidenced in the text
- Include both technical and soft skills
- Look for industry-specific knowledge and domain expertise
- Consider skills implied by job titles, projects, and achievements described

Please respond only with valid JSON in this exact format:
{{
  "technical_skills": ["skill1", "skill2", ...],
  "programming_skills": ["language1", "framework1", ...],
  "leadership_skills": ["skill1", "skill2", ...],
  "collaboration_skills": ["skill1", "skill2", ...],
  "research_skills": ["skill1", "skill2", ...],
  "soft_skills": ["skill1", "skill2", ...],
  "all_skills": ["combined list of all unique skills"]
}}
"""
        
        return prompt
    
    def _call_azure_openai_for_skills(self, prompt: str) -> Optional[str]:
        """Call Azure OpenAI API for skill extraction."""
        try:
            response = self.azure_client.chat.completions.create(
                model=self.azure_openai_deployment,
                messages=[
                    {
                        "role": "system",
                        "content": "You are an expert HR professional and technical recruiter specializing in skill assessment. Analyze documents (resumes, project descriptions, reports, etc.) to extract professional skills and competencies. Focus on both technical skills (programming, tools, platforms) and soft skills (leadership, communication, problem-solving). Be thorough but accurate - only extract skills that are clearly evidenced in the text."
                    },
                    {
                        "role": "user",
                        "content": prompt
                    }
                ],
                max_tokens=2000,
                temperature=0.2,
                top_p=0.9
            )
            
            if response.choices and len(response.choices) > 0:
                return response.choices[0].message.content
            else:
                logger.warning("No response from Azure OpenAI")
                return None
                
        except Exception as e:
            logger.error(f"Error calling Azure OpenAI for skills: {e}")
            return None
    
    def _parse_skills_response(self, ai_response: str) -> Dict[str, List[str]]:
        """Parse AI response and extract categorized skills."""
        try:
            # Try to extract JSON from the response
            json_start = ai_response.find('{')
            json_end = ai_response.rfind('}') + 1
            
            if json_start >= 0 and json_end > json_start:
                json_str = ai_response[json_start:json_end]
                parsed = json.loads(json_str)
                
                # Ensure all required categories exist
                result = {
                    'technical_skills': parsed.get("technical_skills", []),
                    'programming_skills': parsed.get("programming_skills", []),
                    'leadership_skills': parsed.get("leadership_skills", []),
                    'collaboration_skills': parsed.get("collaboration_skills", []),
                    'research_skills': parsed.get("research_skills", []),
                    'soft_skills': parsed.get("soft_skills", []),
                    'all_skills': parsed.get("all_skills", [])
                }
                
                # If all_skills is empty, combine all categories
                if not result['all_skills']:
                    all_skills = []
                    for category in ['technical_skills', 'programming_skills', 'leadership_skills', 
                                   'collaboration_skills', 'research_skills', 'soft_skills']:
                        all_skills.extend(result[category])
                    result['all_skills'] = list(set(all_skills))
                
                logger.info(f"GenAI extracted {len(result['all_skills'])} total skills across {len([k for k, v in result.items() if k != 'all_skills' and v])} categories")
                return result
            else:
                logger.error("Could not find valid JSON in AI response")
                return self._empty_skills_result()
                
        except json.JSONDecodeError as e:
            logger.error(f"JSON parsing error: {e}")
            return self._empty_skills_result()
        except Exception as e:
            logger.error(f"Error parsing skills response: {e}")
            return self._empty_skills_result()
    
    def _fallback_skill_extraction(self, text: str) -> Dict[str, List[str]]:
        """Fallback skill extraction using pattern matching when GenAI is unavailable."""
        logger.info("Using fallback pattern-based skill extraction")
        
        # Simple pattern-based extraction for fallback
        text_lower = text.lower()
        
        # Basic skill patterns for fallback
        technical_patterns = [
            r'\b(?:python|javascript|java|react|typescript|node\.js|sql|aws|azure|docker|git)\b',
            r'\b(?:machine learning|data analysis|api|database|testing|devops|cloud)\b'
        ]
        
        soft_patterns = [
            r'(?:leadership|management|communication|collaboration|problem solving|teamwork)',
            r'(?:research|analysis|strategic planning|project management)'
        ]
        
        technical_skills = []
        soft_skills = []
        
        for pattern in technical_patterns:
            matches = re.findall(pattern, text_lower)
            technical_skills.extend(matches)
        
        for pattern in soft_patterns:
            matches = re.findall(pattern, text_lower)
            soft_skills.extend(matches)
        
        # Clean up and deduplicate
        technical_skills = list(set([skill.title() for skill in technical_skills]))
        soft_skills = list(set([skill.title() for skill in soft_skills]))
        all_skills = technical_skills + soft_skills
        
        return {
            'technical_skills': technical_skills,
            'programming_skills': [s for s in technical_skills if s.lower() in ['python', 'javascript', 'java', 'typescript']],
            'leadership_skills': [s for s in soft_skills if 'leadership' in s.lower() or 'management' in s.lower()],
            'collaboration_skills': [s for s in soft_skills if 'collaboration' in s.lower() or 'teamwork' in s.lower()],
            'research_skills': [s for s in soft_skills if 'research' in s.lower() or 'analysis' in s.lower()],
            'soft_skills': soft_skills,
            'all_skills': all_skills
        }
    
    def _empty_skills_result(self) -> Dict[str, List[str]]:
        """Return empty skills result structure."""
        return {
            'technical_skills': [],
            'programming_skills': [],
            'leadership_skills': [],
            'collaboration_skills': [],
            'research_skills': [],
            'soft_skills': [],
            'all_skills': []
        }
    
    @staticmethod
    def generate_summary(text: str, categorized_skills: Dict[str, List[str]]) -> str:
        """Generate an intelligent summary from document text and categorized skills."""
        word_count = len(text.split())
        
        # Determine document type
        text_lower = text.lower()
        if any(term in text_lower for term in ['resume', 'cv', 'curriculum vitae']):
            doc_type = 'resume'
        elif any(term in text_lower for term in ['slide', 'presentation', 'powerpoint']):
            doc_type = 'presentation'
        elif any(term in text_lower for term in ['project', 'development', 'implementation']):
            doc_type = 'project'
        elif any(term in text_lower for term in ['research', 'study', 'analysis', 'investigation']):
            doc_type = 'research'
        elif any(term in text_lower for term in ['report', 'findings', 'evaluation']):
            doc_type = 'report'
        else:
            doc_type = 'document'
        
        # Generate summary based on skills and content
        summary_parts = []
        
        # Leadership component
        if categorized_skills['leadership_skills']:
            leadership_str = ', '.join(categorized_skills['leadership_skills'][:3])
            summary_parts.append(f"Demonstrates strong leadership capabilities in {leadership_str}")
        
        # Technical component
        if categorized_skills['technical_skills']:
            tech_count = len(categorized_skills['technical_skills'])
            if tech_count > 5:
                summary_parts.append(f"Showcases expertise across {tech_count} technical domains")
            else:
                tech_str = ', '.join(categorized_skills['technical_skills'][:4])
                summary_parts.append(f"Technical proficiency in {tech_str}")
        
        # Programming component
        if categorized_skills['programming_skills']:
            prog_count = len(categorized_skills['programming_skills'])
            if prog_count > 3:
                summary_parts.append(f"Multi-language programming expertise ({prog_count} technologies)")
            else:
                prog_str = ', '.join(categorized_skills['programming_skills'])
                summary_parts.append(f"Programming skills in {prog_str}")
        
        # Research component
        if categorized_skills['research_skills']:
            summary_parts.append("Strong analytical and research capabilities")
        
        # Collaboration component
        if categorized_skills['collaboration_skills']:
            summary_parts.append("Proven collaborative and teamwork abilities")
        
        # Document type specific intro
        if doc_type == 'resume':
            intro = "Professional profile highlighting"
        elif doc_type == 'presentation':
            intro = "Professional presentation covering"
        elif doc_type == 'project':
            intro = "Project documentation showcasing"
        elif doc_type == 'research':
            intro = "Research document demonstrating"
        else:
            intro = "Professional document evidencing"
        
        # Combine summary parts
        if summary_parts:
            main_summary = f"{intro} {', '.join(summary_parts[:3])}."
        else:
            main_summary = f"{intro} professional experience and capabilities."
        
        # Add quantitative information
        if word_count > 2000:
            quantitative = f"This comprehensive {word_count:,}-word document provides detailed insights into advanced professional capabilities."
        elif word_count > 500:
            quantitative = f"This {word_count:,}-word document demonstrates practical experience across multiple domains."
        else:
            quantitative = f"This concise {word_count}-word overview highlights key professional strengths."
        
        return f"{main_summary} {quantitative}"

# Create global DocumentProcessor instance
doc_processor = DocumentProcessor()

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint."""
    return jsonify({"status": "healthy", "service": "document-processor-api"})

@app.route('/process-document', methods=['POST'])
def process_document():
    """Process uploaded document and extract text, skills, and generate summary."""
    try:
        # Validate request
        if 'file' not in request.files:
            return jsonify({"error": "No file provided"}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({"error": "No file selected"}), 400
        
        if not allowed_file(file.filename):
            return jsonify({"error": f"File type not supported. Allowed types: {', '.join(ALLOWED_EXTENSIONS)}"}), 400
        
        start_time = time.time()
        
        # Save file temporarily
        filename = secure_filename(file.filename)
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        file.save(file_path)
        
        try:
            # Determine file type and process accordingly
            file_extension = filename.lower().split('.')[-1]
            
            if file_extension == 'pdf':
                text, metadata = DocumentProcessor.extract_pdf_text_advanced(file_path)
            elif file_extension == 'docx':
                text, metadata = DocumentProcessor.extract_docx_text(file_path)
            elif file_extension == 'pptx':
                text, metadata = DocumentProcessor.extract_pptx_text(file_path)
            else:
                return jsonify({"error": "Unsupported file type"}), 400
            
            # Extract skills and generate summary using GenAI
            categorized_skills = doc_processor.extract_skills_with_genai(text)
            summary = DocumentProcessor.generate_summary(text, categorized_skills)
            
            # Calculate processing time
            processing_time = int((time.time() - start_time) * 1000)
            
            # Build response
            response_data = {
                "text": text,
                "skills": categorized_skills['all_skills'],
                "categorized_skills": categorized_skills,
                "summary": summary,
                "metadata": {
                    **metadata,
                    "wordCount": len(text.split()),
                    "fileType": file_extension.upper(),
                    "processingTime": processing_time,
                    "fileName": filename,
                    "skillCategories": {
                        "technical": len(categorized_skills['technical_skills']),
                        "programming": len(categorized_skills['programming_skills']),
                        "leadership": len(categorized_skills['leadership_skills']),
                        "collaboration": len(categorized_skills['collaboration_skills']),
                        "research": len(categorized_skills['research_skills']),
                        "soft": len(categorized_skills['soft_skills'])
                    }
                }
            }
            
            logger.info(f"Successfully processed {filename}: {len(text)} chars, {len(categorized_skills['all_skills'])} skills")
            return jsonify(response_data)
            
        finally:
            # Clean up temporary file
            if os.path.exists(file_path):
                os.remove(file_path)
    
    except Exception as e:
        logger.error(f"Error processing document: {str(e)}")
        return jsonify({"error": f"Processing failed: {str(e)}"}), 500

@app.route('/extract-text', methods=['POST'])
def extract_text_only():
    """Extract only text from document (simpler endpoint)."""
    try:
        if 'file' not in request.files:
            return jsonify({"error": "No file provided"}), 400
        
        file = request.files['file']
        if not allowed_file(file.filename):
            return jsonify({"error": "File type not supported"}), 400
        
        filename = secure_filename(file.filename)
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        file.save(file_path)
        
        try:
            file_extension = filename.lower().split('.')[-1]
            
            if file_extension == 'pdf':
                text, metadata = DocumentProcessor.extract_pdf_text_advanced(file_path)
            elif file_extension == 'docx':
                text, metadata = DocumentProcessor.extract_docx_text(file_path)
            elif file_extension == 'pptx':
                text, metadata = DocumentProcessor.extract_pptx_text(file_path)
            
            return jsonify({
                "text": text,
                "metadata": metadata
            })
            
        finally:
            if os.path.exists(file_path):
                os.remove(file_path)
    
    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == '__main__':
    logger.info("Starting Document Processing API with Azure OpenAI integration...")
    logger.info(f"Allowed file types: {ALLOWED_EXTENSIONS}")
    logger.info(f"Max file size: {app.config['MAX_CONTENT_LENGTH'] / (1024*1024):.1f}MB")
    
    # Check Azure OpenAI configuration
    if doc_processor.azure_client:
        logger.info("✅ Azure OpenAI client is ready for GenAI skill extraction")
    else:
        logger.warning("⚠️ Azure OpenAI client not available - using fallback pattern-based extraction")
    
    app.run(debug=True, host='0.0.0.0', port=5001)
